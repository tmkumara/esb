"""
FineXaTech ESB Platform — Architecture & Demo PowerPoint Generator
Run: python generate_pptx.py
Output: ESB_Platform_Demo.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────
NAVY    = RGBColor(0x1e, 0x3a, 0x8a)   # headings
BLUE    = RGBColor(0x2b, 0x6c, 0xd8)   # accent
LAVENDER= RGBColor(0xf0, 0xf4, 0xff)   # slide bg
WHITE   = RGBColor(0xff, 0xff, 0xff)
ORANGE  = RGBColor(0xf9, 0x73, 0x16)
GREEN   = RGBColor(0x16, 0xa3, 0x4a)
RED     = RGBColor(0xdc, 0x26, 0x26)
GRAY    = RGBColor(0x94, 0xa3, 0xb8)
DARK    = RGBColor(0x1e, 0x29, 0x3b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # totally blank

# ─── helpers ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def bg(slide, color=LAVENDER):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.2, fill=NAVY)
    add_text(slide, title, 0.35, 0.1, 10, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.75, 10, 0.4, size=13, color=RGBColor(0xbf,0xdb,0xff))

def card(slide, l, t, w, h, title="", body_lines=None, title_color=NAVY,
         border=BLUE, bg_color=WHITE):
    add_rect(slide, l, t, w, h, fill=bg_color, line=border, line_w=Pt(1.2))
    if title:
        add_rect(slide, l, t, w, 0.38, fill=border)
        add_text(slide, title, l+0.12, t+0.04, w-0.2, 0.32,
                 size=11, bold=True, color=WHITE)
    if body_lines:
        body = "\n".join(body_lines)
        add_text(slide, body, l+0.15, t+0.45, w-0.25, h-0.55,
                 size=10, color=DARK)

def arrow(slide, x1, y, x2, color=BLUE, label=""):
    """Horizontal arrow from x1→x2 at height y (in inches)."""
    from pptx.util import Emu, Pt
    connector = slide.shapes.add_connector(1,
        Inches(x1), Inches(y), Inches(x2), Inches(y))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    if label:
        add_text(slide, label, (x1+x2)/2-0.4, y-0.28, 0.8, 0.25,
                 size=8, color=color, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)

# gradient-like bg
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 4.2, 13.33, 3.3, fill=RGBColor(0x17, 0x2b, 0x6a))

# decorative bar
add_rect(s, 0, 3.9, 13.33, 0.06, fill=BLUE)

add_text(s, "FineXaTech ESB Platform",
         1, 1.0, 11, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Enterprise Service Bus — Architecture & Live Demo",
         1, 2.3, 11, 0.6, size=20, color=RGBColor(0xbf,0xdb,0xff), align=PP_ALIGN.CENTER)

# tech pills
pills = [("Apache Camel 4.7", 2.2), ("Spring Boot 3.3", 4.5), ("React 18 + Vite", 6.8)]
for label, lx in pills:
    add_rect(s, lx, 3.25, 2.0, 0.45, fill=BLUE, line=WHITE, line_w=Pt(0.5))
    add_text(s, label, lx+0.1, 3.28, 1.85, 0.4, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "Phase 1 Complete  ·  Fully Operational",
         1, 4.6, 11, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, "FineXaTech Engineering  ·  2026",
         1, 6.8, 11, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Platform Overview
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Platform Overview", "What FineXaTech ESB does — in one picture")

# big flow diagram boxes
boxes = [
    (0.4, "Client\n(Mobile / API)", BLUE),
    (3.1, "ESB Runtime\n:9090", NAVY),
    (5.9, "Transform\nEngine", RGBColor(0x6d,0x28,0xd9)),
    (8.7, "Backend\nServices", GREEN),
]
for lx, lbl, col in boxes:
    add_rect(s, lx, 2.4, 2.3, 1.5, fill=col)
    add_text(s, lbl, lx+0.1, 2.55, 2.1, 1.2, size=13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

# arrows
for x in [2.7, 5.4, 8.2]:
    arrow(s, x, 3.15, x+0.4, color=ORANGE)

# labels below boxes
labels = ["REST/HTTP", "Route\nAssembler", "Jolt · XSLT\nGroovy", "SOAP · REST\nJMS · Kafka"]
for i, (lx, _, _) in enumerate(boxes):
    add_text(s, labels[i], lx, 4.05, 2.3, 0.5, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# UI bar
add_rect(s, 0.4, 5.3, 12.5, 0.9, fill=RGBColor(0xe0,0xe7,0xff), line=BLUE, line_w=Pt(1))
add_text(s, "ESB UI (port 3000)  —  Routes · Builder · Validation · Monitoring  —  React + Vite + Tailwind",
         1, 5.45, 11.5, 0.6, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# key callouts
callouts = [
    (0.5, 1.5, "Low-Code / No-Code\nroute building"),
    (4.0, 1.5, "5-Layer Validation\npipeline"),
    (7.5, 1.5, "Hot Reload\nno restart needed"),
    (10.5, 1.5, "Mock Server\nbuilt-in"),
]
for lx, ly, txt in callouts:
    add_rect(s, lx, ly, 2.3, 0.75, fill=WHITE, line=BLUE, line_w=Pt(0.8))
    add_text(s, txt, lx+0.1, ly+0.05, 2.1, 0.65, size=10, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Module Structure
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Module Structure", "Maven multi-module layout")

modules = [
    ("esb-spec", GRAY,
     ["Pure Java POJOs", "RouteSpec, TargetSpec", "TransformSpec, SourceSpec", "No Spring dependency"]),
    ("esb-compiler", RGBColor(0x6d,0x28,0xd9),
     ["ValidationPipeline (5 layers)", "RouteAssembler", "CamelDryRunCompiler", "SpecRule interface"]),
    ("esb-adapters", BLUE,
     ["RestSourceAdapter", "SoapTargetAdapter", "GroovyTransformAdapter", "MockResponseTargetAdapter"]),
    ("esb-interceptors", ORANGE,
     ["CorrelationInterceptor", "RetryInterceptor", "AuthInterceptor", "MetricsInterceptor"]),
    ("esb-runtime", NAVY,
     ["Spring Boot :9090", "LiveRouteRegistry", "HotReloadWatcher", "Management API /manage/**"]),
    ("esb-ui", GREEN,
     ["React 18 + Vite", "Route Builder (drag-drop)", "Validation page", "Monitoring page"]),
]

for i, (name, col, lines) in enumerate(modules):
    col_idx = i % 3
    row_idx = i // 3
    lx = 0.35 + col_idx * 4.3
    ly = 1.45 + row_idx * 2.7
    card(s, lx, ly, 4.0, 2.4, title=name, body_lines=lines,
         border=col, bg_color=WHITE)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Validation Pipeline
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "5-Layer Validation Pipeline", "Every route passes all 5 gates before going live")

layers = [
    ("L1\nSTRUCTURAL", "Required fields\nYAML format", BLUE),
    ("L2\nSCHEMA", "Enum values\nAllowed types", RGBColor(0x6d,0x28,0xd9)),
    ("L3\nSEMANTIC", "Logic checks\nFile/URL exists", ORANGE),
    ("L4\nCOMPATIBILITY", "Source ↔ Target\ntype match", RGBColor(0x0e,0x7a,0x70)),
    ("L5\nDRY_RUN", "Real Camel context\nmock endpoints", GREEN),
]

for i, (lbl, desc, col) in enumerate(layers):
    lx = 0.5 + i * 2.45
    add_rect(s, lx, 1.7, 2.1, 2.5, fill=col)
    add_text(s, lbl, lx+0.1, 1.85, 1.9, 0.9, size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, lx+0.1, 2.85, 1.9, 1.0, size=11,
             color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, lx+2.1, 2.95, lx+2.45, color=NAVY, label="")

add_text(s, "✓  Stops at first FAIL — no partial deploys",
         1.5, 4.55, 10, 0.5, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_text(s,
    "RequiredFieldsRule · HttpMethodRule · EnvVarResolvableRule · CompatibilityMatrix · CamelDryRunCompiler",
    1.0, 5.15, 11, 0.4, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 5 — REST → SOAP Demo Flow
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Live Demo — REST → SOAP Account Balance", "account-balance route in action")

steps = [
    (0.4,  "Mobile Client",       "GET /api/v1/accounts\n/12345/balance\n(REST, no body)",        BLUE),
    (3.2,  "Groovy\nReq Transform","headers['accountId']\n→ SOAP XML Envelope",                    RGBColor(0x6d,0x28,0xd9)),
    (6.0,  "SOAP Target",         "POST /mock/soap/\nbalance-service\ntext/xml",                   NAVY),
    (8.8,  "Mock Response\nTarget","Returns SOAP XML\n(no Java code)\nconfigured in YAML",         ORANGE),
]

for i, (lx, title, desc, col) in enumerate(steps):
    add_rect(s, lx, 1.7, 2.3, 2.2, fill=col)
    add_text(s, title, lx+0.1, 1.8, 2.1, 0.55, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, desc, lx+0.1, 2.4, 2.1, 1.3, size=10,
             color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, lx+2.3, 2.8, lx+2.65, color=ORANGE)

# Response path
resp_steps = [
    (8.8, "SOAP XML\nResponse", GREEN),
    (6.0, "SoapTargetAdapter\nXML → JSON (auto)", NAVY),
    (3.2, "Jolt\nRes Transform", RGBColor(0x6d,0x28,0xd9)),
    (0.4, "Clean JSON\nto Mobile Client", GREEN),
]
for i, (lx, title, col) in enumerate(resp_steps):
    add_rect(s, lx, 4.3, 2.3, 1.0, fill=col)
    add_text(s, title, lx+0.1, 4.4, 2.1, 0.8, size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(s, lx, 4.8, lx-0.35, color=GREEN)

add_text(s, "↑ Request flow", 11.2, 2.6, 1.8, 0.3, size=9, color=ORANGE)
add_text(s, "↓ Response flow", 11.2, 4.6, 1.8, 0.3, size=9, color=GREEN)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Adapter Types
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Adapter Catalogue", "Current Phase 1 implementations")

adapters = [
    ("SOURCE ADAPTERS", BLUE, [
        ("rest",   "REST/HTTP",     "Any REST endpoint — GET/POST/PUT/DELETE"),
        ("direct", "Direct",        "In-process Camel direct channel"),
        ("jms",    "JMS (Phase 2)", "Queue / topic source"),
        ("kafka",  "Kafka (Phase 2)","Consumer from Kafka topic"),
    ]),
    ("TARGET ADAPTERS", NAVY, [
        ("soap",          "SOAP Target",    "HTTP POST with SOAPAction + XML→JSON auto-convert"),
        ("http",          "HTTP Target",    "REST outbound call"),
        ("mock-response", "Mock Response",  "Static YAML-configured response — zero Java code"),
        ("mock-echo",     "Mock Echo",      "Returns Groovy-computed body as HTTP 200"),
    ]),
    ("TRANSFORM ADAPTERS", RGBColor(0x6d,0x28,0xd9), [
        ("jolt",        "Jolt JSON",    "Declarative JSON→JSON with Field Mapper UI"),
        ("xslt",        "XSLT",         "XML→XML stylesheet transform"),
        ("groovy",      "Groovy Script", "Any→any scripted transform with body+headers access"),
        ("passthrough", "Passthrough",   "No transform — body passes unchanged"),
    ]),
]

col_offsets = [0.35, 4.55, 8.75]
for ci, (section, col, items) in enumerate(adapters):
    lx = col_offsets[ci]
    add_rect(s, lx, 1.45, 4.1, 0.38, fill=col)
    add_text(s, section, lx+0.1, 1.5, 3.9, 0.3, size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    for ri, (key, name, desc) in enumerate(items):
        ly = 1.95 + ri * 1.25
        add_rect(s, lx, ly, 4.1, 1.1, fill=WHITE, line=col, line_w=Pt(0.8))
        add_rect(s, lx, ly, 1.0, 1.1, fill=RGBColor(0xee,0xf2,0xff))
        add_text(s, key, lx+0.05, ly+0.25, 0.95, 0.6, size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, name, lx+1.1, ly+0.08, 3.0, 0.35, size=11, bold=True, color=DARK)
        add_text(s, desc, lx+1.1, ly+0.48, 2.9, 0.5, size=9, color=GRAY)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Issues Found & Fixed
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Issues Found & Fixed", "Key bugs resolved during Phase 1 development")

issues = [
    ("YAML TAB error",
     "Groovy inline block had TAB characters",
     "Use | block literal with spaces. Single-line GString avoids it entirely."),
    ("Groovy 'No such property'",
     "def id creates local var; ${id} in GString uses binding lookup",
     "Inline expression directly: ${headers['accountId'] ?: 'UNKNOWN'}"),
    ("Path doubling — 404",
     "Camel HTTP appends source path to target URL even with bridgeEndpoint=true",
     "Remove HTTP_PATH, HTTP_URI, HTTP_QUERY headers in SoapTargetAdapter.preProcessor()"),
    ("Transforms showing 'none'",
     "RouteStatusView record had no transform fields",
     "Added requestTransformType + responseTransformType to RouteStatusView"),
    ("ValidationPage crash",
     "Backend returns {messages[]} but UI expects {layers[]}",
     "Added normalizeValidationResponse() to group messages by layer"),
    ("endpointUrl required for mock",
     "RequiredFieldsRule always checked endpointUrl for all targets",
     "Skip check when target.type = mock-response or mock-echo"),
    ("Groovy preview 'headers' error",
     "TransformPreviewService only injected body, not headers",
     "Added headers map + Sample Headers JSON input in CodeEditorModal"),
]

for i, (title, cause, fix) in enumerate(issues):
    row = i % 4
    col = i // 4
    lx = 0.35 + col * 6.5
    ly = 1.45 + row * 1.45
    add_rect(s, lx, ly, 6.2, 1.35, fill=WHITE, line=BLUE, line_w=Pt(0.8))
    add_rect(s, lx, ly, 6.2, 0.33, fill=NAVY)
    add_text(s, title, lx+0.1, ly+0.03, 6.0, 0.28, size=10, bold=True, color=WHITE)
    add_text(s, f"⚠ {cause}", lx+0.1, ly+0.38, 6.0, 0.38, size=9, color=RED)
    add_text(s, f"✓ {fix}",   lx+0.1, ly+0.76, 6.0, 0.48, size=9, color=GREEN)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Route Builder UI
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Route Builder — Low Code UI", "Drag-drop canvas → auto-generates YAML")

# Canvas area
add_rect(s, 0.35, 1.4, 8.5, 5.5, fill=WHITE, line=BLUE, line_w=Pt(1))
add_text(s, "Canvas (ReactFlow @xyflow/react v12)", 0.5, 1.45, 6, 0.3, size=9, color=GRAY)

# Nodes on canvas
node_types = [
    (1.0, 2.5, "SOURCE", "REST GET\n/v1/accounts/{id}", BLUE),
    (2.8, 2.5, "REQ\nTRANSFORM", "groovy\nScript", RGBColor(0x6d,0x28,0xd9)),
    (4.6, 2.5, "TARGET", "SOAP\nbalance-service", NAVY),
    (2.8, 4.5, "RES\nTRANSFORM", "jolt\nFlatten", RGBColor(0x6d,0x28,0xd9)),
]
for lx, ly, label, sub, col in node_types:
    add_rect(s, lx, ly, 1.6, 1.3, fill=col)
    add_text(s, label, lx+0.05, ly+0.08, 1.5, 0.45, size=9, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sub, lx+0.05, ly+0.55, 1.5, 0.65, size=9,
             color=WHITE, align=PP_ALIGN.CENTER)

# Right panel
add_rect(s, 9.0, 1.4, 4.0, 5.5, fill=RGBColor(0xf8,0xfa,0xff), line=GRAY, line_w=Pt(0.5))
add_text(s, "PROPERTIES", 9.1, 1.5, 3.8, 0.3, size=10, bold=True, color=NAVY)
props = [
    ("Label", "SOAP Target"),
    ("Endpoint URL", "http://localhost:9090/mock/..."),
    ("Timeout (ms)", "5000"),
    ("", ""),
    ("YAML Preview", ""),
    ("target:", ""),
    ("  type: soap", ""),
    ("  endpointUrl: ...", ""),
    ("  timeout:", ""),
    ("    readMs: 5000", ""),
]
for i, (k, v) in enumerate(props):
    ly = 1.95 + i * 0.44
    if k:
        add_text(s, k, 9.1, ly, 1.5, 0.38, size=9, bold=True, color=DARK)
    if v:
        add_text(s, v, 10.6, ly, 2.3, 0.38, size=9, color=GRAY)

# toolbar
toolbar = ["Reset", "YAML", "Validate", "Deploy", "Save to Disk"]
colors  = [GRAY, GRAY, BLUE, GREEN, NAVY]
for i, (lbl, col) in enumerate(zip(toolbar, colors)):
    add_rect(s, 0.5 + i*1.65, 6.55, 1.45, 0.4, fill=col)
    add_text(s, lbl, 0.5 + i*1.65, 6.58, 1.45, 0.34,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Mock Server Concept
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Built-in Mock Server", "Zero Java code — configure mocks as ESB routes")

# Before / After
add_text(s, "BEFORE  (Java mock classes)", 0.5, 1.45, 5.8, 0.4, size=14, bold=True, color=RED)
add_rect(s, 0.35, 1.9, 5.9, 4.5, fill=WHITE, line=RED, line_w=Pt(1))
before = """\
// New Java class needed per mock
@RestController
@Profile("demo")
public class MockSoapController {

  @PostMapping("/mock/soap/balance-service")
  public String balance(...) {
    return "<soap:Envelope>...</soap:Envelope>";
  }
  // Repeat for every new mock endpoint
  // Requires: code change + rebuild + restart
}"""
add_text(s, before, 0.5, 2.0, 5.7, 4.2, size=10, color=RED)

add_text(s, "AFTER  (mock-response target)", 7.0, 1.45, 5.8, 0.4, size=14, bold=True, color=GREEN)
add_rect(s, 6.85, 1.9, 6.0, 4.5, fill=WHITE, line=GREEN, line_w=Pt(1))
after = """\
# New mock = new route in the UI
# No Java. No restart. No rebuild.

target:
  type: mock-response
  mockStatusCode: 200
  mockBody: |
    <?xml version="1.0"?>
    <soap:Envelope ...>
      <soap:Body>
        <GetAccountBalanceResponse>
          <balance>2500.75</balance>
        </GetAccountBalanceResponse>
      </soap:Body>
    </soap:Envelope>"""
add_text(s, after, 7.0, 2.0, 5.7, 4.2, size=10, color=GREEN)

add_text(s, "→", 6.3, 4.0, 0.6, 0.5, size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 10 — Current Status & Roadmap
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
header_bar(s, "Current Status & Roadmap", "Phase 1 complete — what's next")

phases = [
    ("Phase 1\n✅ Complete", NAVY, [
        "REST source adapter",
        "SOAP + HTTP target adapters",
        "Jolt · XSLT · Groovy transforms",
        "5-layer validation pipeline",
        "Hot-reload file watcher",
        "Full React UI (Builder, Routes, Validation)",
        "Mock Response target (zero-code mocks)",
        "Groovy preview with headers",
        "Response transform + Jolt flatten",
    ]),
    ("Phase 2\n🔄 Planned", BLUE, [
        "JMS source + target",
        "Kafka source + target",
        "File + SFTP source/target",
        "Timer source",
        "Multi-broker orchestration UI",
        "Route versioning + rollback",
        "Auth interceptor (JWT/OAuth2)",
        "Metrics dashboard (Prometheus)",
    ]),
    ("Phase 3\n📋 Demand-driven", GRAY, [
        "FIXML adapter",
        "JDBC / database target",
        "gRPC source + target",
        "SAP RFC adapter",
        "Visual XSLT editor",
        "Route templates library",
        "Git-based route deployment",
        "Multi-tenant isolation",
    ]),
]

for ci, (title, col, items) in enumerate(phases):
    lx = 0.35 + ci * 4.3
    add_rect(s, lx, 1.45, 4.1, 0.65, fill=col)
    add_text(s, title, lx+0.1, 1.5, 3.9, 0.55, size=14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    for ri, item in enumerate(items):
        ly = 2.22 + ri * 0.56
        mark = "●" if col != GRAY else "○"
        add_text(s, f"{mark}  {item}", lx+0.2, ly, 3.8, 0.5, size=10,
                 color=DARK if col != GRAY else GRAY)

# ══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Thank You / Questions
# ══════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 5.5, 13.33, 2.0, fill=RGBColor(0x17,0x2b,0x6a))
add_rect(s, 0, 5.38, 13.33, 0.06, fill=BLUE)

add_text(s, "Thank You", 1, 1.2, 11, 1.2, size=54, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "FineXaTech ESB Platform — Phase 1",
         1, 2.7, 11, 0.6, size=20, color=RGBColor(0xbf,0xdb,0xff), align=PP_ALIGN.CENTER)

links = [
    ("Backend", "http://localhost:9090/manage/health"),
    ("UI",      "http://localhost:3000"),
    ("Routes",  "http://localhost:9090/manage/routes"),
]
for i, (label, url) in enumerate(links):
    lx = 1.5 + i * 3.5
    add_rect(s, lx, 3.7, 3.0, 0.8, fill=BLUE)
    add_text(s, label, lx+0.1, 3.75, 2.8, 0.3, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, url,   lx+0.1, 4.07, 2.8, 0.35, size=9, color=RGBColor(0xbf,0xdb,0xff), align=PP_ALIGN.CENTER)

add_text(s, "Questions?",
         1, 5.6, 11, 0.6, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────
out = r"D:\FineXaTech\POC\esb\docs\ESB_Platform_Demo.pptx"
prs.save(out)
print(f"Saved: {out}")
